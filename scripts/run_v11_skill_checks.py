#!/usr/bin/env python3
"""Run the gen-tgo-ppt-skill V1.1 maintenance harness."""

from __future__ import annotations

import argparse
import ast
import copy
import hashlib
import importlib.util
import json
import os
import re
import subprocess
import sys
import tempfile
from zipfile import ZipFile
from collections import Counter
from contextlib import contextmanager
from pathlib import Path
from typing import Any, Callable, Iterator


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
COMPATIBILITY_WRAPPER = SCRIPTS / "run_v1_skill_checks.py"
SKILL_MANIFEST_PATH = ROOT / "references" / "skill-manifest.json"
TEMPLATE_MANIFEST_PATH = ROOT / "references" / "template-manifest.json"
EXPECTED_SKILL_VERSION = "V1.1"
EXPECTED_TEMPLATE_VERSION = "V1.1"
EXPECTED_LAYOUT_VERSION = "V1"
EXPECTED_SCHEMA = "gen-tgo-ppt-skill-manifest-v1.1"
EVAL_SCHEMA = "gen-tgo-ppt-routing-eval-v1.1"
ACTIVATION_SAFETY_SCHEMA = "gen-tgo-ppt-activation-safety-v1.1"
PROTECTED_SHA256 = "e1bbdc853c5d013810730f118850ad46dda9ab3bf59d813a4db890e283757c18"
ENTRYPOINT_ENV = "GEN_TGO_PPT_HARNESS_ENTRYPOINT"
TEXT_EXTENSIONS = {".json", ".md", ".py", ".txt", ".yaml", ".yml"}
CASE_TYPES = {"create", "convert", "repair", "check_only", "no_trigger", "handoff"}
IMPORT_TO_DISTRIBUTION = {"pptx": "python-pptx", "lxml": "lxml", "PIL": "Pillow"}
TEMP_ROOT_ENV = "GEN_TGO_PPT_TEMP_ROOT"
TEMP_ROOT: Path | None = None
TEMP_REQUIRED_CHECKS = {
    "generation_log_dry_run_and_actual",
    "html_fixtures",
    "pptx_fixtures",
    "helper_contracts",
    "evaluator_fixtures_and_negative_paths",
}


def resolve_temp_root(explicit: Path | None) -> tuple[Path | None, dict[str, Any]]:
    source = "--temp-root" if explicit is not None else TEMP_ROOT_ENV if os.environ.get(TEMP_ROOT_ENV) else "system"
    candidate = explicit
    if candidate is None and os.environ.get(TEMP_ROOT_ENV):
        candidate = Path(os.environ[TEMP_ROOT_ENV])
    if candidate is None:
        try:
            candidate = Path(tempfile.gettempdir())
        except (FileNotFoundError, OSError) as exc:
            return None, {
                "status": "INSUFFICIENT",
                "source": source,
                "error_code": "E_TEMP_ROOT_UNAVAILABLE",
                "error": f"{type(exc).__name__}: {exc}",
            }

    candidate = candidate.expanduser().resolve()
    if not candidate.is_dir():
        return None, {
            "status": "INSUFFICIENT",
            "source": source,
            "path": str(candidate),
            "error_code": "E_TEMP_ROOT_UNAVAILABLE",
            "error": "configured temporary root must already exist and be a directory",
        }
    try:
        with tempfile.TemporaryDirectory(prefix="gen-tgo-ppt-probe-", dir=candidate) as probe:
            probe_file = Path(probe) / "write-probe"
            probe_file.write_text("ok", encoding="utf-8")
            if probe_file.read_text(encoding="utf-8") != "ok":
                raise OSError("temporary workspace write/read probe mismatch")
    except (FileNotFoundError, OSError, PermissionError) as exc:
        return None, {
            "status": "INSUFFICIENT",
            "source": source,
            "path": str(candidate),
            "error_code": "E_TEMP_ROOT_UNAVAILABLE",
            "error": f"{type(exc).__name__}: {exc}",
        }
    return candidate, {"status": "PASS", "source": source, "path": str(candidate)}


@contextmanager
def temporary_directory() -> Iterator[str]:
    if TEMP_ROOT is None:
        raise RuntimeError("E_TEMP_ROOT_UNAVAILABLE: deterministic fixture checks require a writable temp root")
    with tempfile.TemporaryDirectory(prefix="gen-tgo-ppt-harness-", dir=TEMP_ROOT) as temporary:
        yield temporary


def run(
    command: list[str],
    *,
    cwd: Path | None = None,
    expected: int | set[int] | None = 0,
) -> subprocess.CompletedProcess[str]:
    env = os.environ.copy()
    env["PYTHONDONTWRITEBYTECODE"] = "1"
    result = subprocess.run(command, cwd=cwd, env=env, text=True, capture_output=True)
    accepted = None if expected is None else ({expected} if isinstance(expected, int) else expected)
    if accepted is not None and result.returncode not in accepted:
        raise AssertionError(
            "Command returned an unexpected exit code:\n"
            f"command={' '.join(command)}\nreturncode={result.returncode}\n"
            f"stdout={result.stdout[-4000:]}\nstderr={result.stderr[-4000:]}"
        )
    return result


def load_json(path: Path) -> dict[str, Any]:
    try:
        value = json.loads(path.read_text(encoding="utf-8"))
    except FileNotFoundError as exc:
        raise AssertionError(f"Required JSON file is missing: {path.relative_to(ROOT)}") from exc
    except json.JSONDecodeError as exc:
        raise AssertionError(
            f"Invalid JSON in {path.relative_to(ROOT)} at line {exc.lineno}, column {exc.colno}"
        ) from exc
    if not isinstance(value, dict):
        raise AssertionError(f"JSON root must be an object: {path.relative_to(ROOT)}")
    return value


def parse_json_stdout(result: subprocess.CompletedProcess[str], label: str) -> Any:
    try:
        return json.loads(result.stdout)
    except json.JSONDecodeError as exc:
        raise AssertionError(f"{label} did not emit valid JSON: {result.stdout[-2000:]}") from exc


def iter_string_values(value: Any):
    if isinstance(value, dict):
        for child in value.values():
            yield from iter_string_values(child)
    elif isinstance(value, list):
        for child in value:
            yield from iter_string_values(child)
    elif isinstance(value, str):
        yield value


def declared_paths(value: Any) -> set[str]:
    prefixes = ("assets/", "references/", "scripts/", "agents/")
    paths = {item for item in iter_string_values(value) if item.startswith(prefixes)}
    for item in iter_string_values(value):
        if item in {"SKILL.md", "README.md"}:
            paths.add(item)
    return paths


def assert_paths_exist(paths: set[str], label: str) -> None:
    missing = sorted(path for path in paths if not (ROOT / path).exists())
    if missing:
        raise AssertionError(f"{label} declares missing paths: {missing}")


def check_frontmatter_and_protected_section() -> dict[str, Any]:
    raw = (ROOT / "SKILL.md").read_bytes()
    frontmatter_match = re.match(rb"\A---\n(.*?)\n---\n", raw, flags=re.S)
    if not frontmatter_match:
        raise AssertionError("SKILL.md must start with YAML frontmatter")
    fields: dict[str, str] = {}
    for line in frontmatter_match.group(1).decode("utf-8").splitlines():
        if ":" not in line:
            raise AssertionError(f"Invalid SKILL.md frontmatter line: {line!r}")
        key, value = line.split(":", 1)
        key = key.strip()
        if key in fields:
            raise AssertionError(f"Duplicate SKILL.md frontmatter field: {key}")
        fields[key] = value.strip()
    if set(fields) != {"name", "description"}:
        raise AssertionError(f"Frontmatter fields must be exactly name + description, got {sorted(fields)}")
    if fields["name"] != "gen-tgo-ppt-skill" or not fields["description"]:
        raise AssertionError("SKILL.md frontmatter name/description values are invalid")

    start_marker = "## 重要说明和介绍".encode("utf-8")
    end_marker = "## 版本".encode("utf-8")
    if raw.count(start_marker) != 1 or raw.count(end_marker) != 1:
        raise AssertionError("Protected section anchors must each occur exactly once")
    start = raw.index(start_marker)
    end = raw.index(end_marker, start)
    between = raw[start:end]
    if not between.endswith(b"\n\n"):
        raise AssertionError("Protected section separator changed")
    protected_bytes = between[:-1]
    digest = hashlib.sha256(protected_bytes).hexdigest()
    if digest != PROTECTED_SHA256:
        raise AssertionError(
            "E_GOVERNANCE_PROTECTED_SECTION_CHANGED: protected SKILL.md bytes changed; "
            f"expected={PROTECTED_SHA256} actual={digest}"
        )
    headings = re.findall(rb"(?m)^## [^\n]+", raw)
    if headings[:2] != [start_marker, end_marker]:
        raise AssertionError("Protected section must remain the first H2 and immediately precede the version section")
    return {"frontmatter_fields": sorted(fields), "protected_sha256": digest}


def check_version_manifests_agents_and_paths() -> dict[str, Any]:
    skill_manifest = load_json(SKILL_MANIFEST_PATH)
    template_manifest = load_json(TEMPLATE_MANIFEST_PATH)
    if skill_manifest.get("schema_version") != EXPECTED_SCHEMA:
        raise AssertionError("Skill manifest schema_version mismatch")
    skill = skill_manifest.get("skill")
    if not isinstance(skill, dict) or skill.get("name") != "gen-tgo-ppt-skill":
        raise AssertionError("Skill manifest identity mismatch")
    if skill.get("version") != EXPECTED_SKILL_VERSION:
        raise AssertionError("Skill manifest version must be V1.1")
    ownership = skill_manifest.get("version_ownership")
    expected_ownership = {
        "skill_contract": EXPECTED_SKILL_VERSION,
        "template_bundle": EXPECTED_TEMPLATE_VERSION,
        "layout_safety_algorithm": EXPECTED_LAYOUT_VERSION,
    }
    if not isinstance(ownership, dict) or any(ownership.get(k) != v for k, v in expected_ownership.items()):
        raise AssertionError("Skill/template/layout version ownership is inconsistent")
    protected = skill_manifest.get("protected_sections")
    if not isinstance(protected, list) or len(protected) != 1:
        raise AssertionError("Skill manifest must declare exactly one protected section")
    if protected[0].get("file") != "SKILL.md" or protected[0].get("sha256") != PROTECTED_SHA256:
        raise AssertionError("Skill manifest protected-section contract mismatch")

    entrypoints = skill_manifest.get("entrypoints")
    if not isinstance(entrypoints, dict):
        raise AssertionError("Skill manifest entrypoints must be an object")
    expected_entrypoints = {
        "current_harness": "scripts/run_v11_skill_checks.py",
        "compatibility_harness": "scripts/run_v1_skill_checks.py",
        "compatibility_direction": "run_v1_skill_checks.py -> run_v11_skill_checks.py",
    }
    if any(entrypoints.get(key) != value for key, value in expected_entrypoints.items()):
        raise AssertionError("Harness entrypoint direction is inconsistent")
    if "run_v11_skill_checks.py -> run_v1_skill_checks.py" not in entrypoints.get("forbidden_calls", []):
        raise AssertionError("Skill manifest must forbid V1.1 -> V1 harness callbacks")

    if template_manifest.get("manifest_type") != "template_bundle":
        raise AssertionError("template-manifest.json manifest_type must be template_bundle")
    if template_manifest.get("version") != "V1":
        raise AssertionError("template-manifest.json legacy version must remain V1")
    if template_manifest.get("template_bundle_version") != EXPECTED_TEMPLATE_VERSION:
        raise AssertionError("template bundle version must be V1.1")
    if template_manifest.get("skill_manifest") != "references/skill-manifest.json":
        raise AssertionError("Template manifest must link to the Skill manifest")

    skill_text = (ROOT / "SKILL.md").read_text(encoding="utf-8")
    readme_text = (ROOT / "README.md").read_text(encoding="utf-8")
    agents_text = (ROOT / "agents" / "openai.yaml").read_text(encoding="utf-8")
    log_source = (SCRIPTS / "create_generation_log.py").read_text(encoding="utf-8")
    if "当前版本：V1.1" not in skill_text or "模板包版本：V1.1" not in skill_text:
        raise AssertionError("SKILL.md version declarations are inconsistent")
    if "Skill 版本：V1.1" not in readme_text or "run_v11_skill_checks.py" not in readme_text:
        raise AssertionError("README.md must identify V1.1 and the canonical harness")
    if (
        "V1.1" not in agents_text
        or "allow_implicit_invocation: true" not in agents_text
        or "E_ROUTE_AMBIGUOUS" not in agents_text
    ):
        raise AssertionError("agents/openai.yaml V1.1 routing metadata is incomplete")
    if (
        'SKILL_VERSION = "V1.1"' not in log_source
        or 'LAYOUT_SAFETY_VERSION = "V1"' not in log_source
        or 'TEMPLATE_BUNDLE_VERSION = "V1.1"' not in log_source
    ):
        raise AssertionError("Generation-log version constants are inconsistent")

    manifest_script_paths = {
        item.get("path") for item in skill_manifest.get("scripts", []) if isinstance(item, dict)
    }
    actual_script_paths = {path.relative_to(ROOT).as_posix() for path in SCRIPTS.glob("*.py")}
    if manifest_script_paths != actual_script_paths:
        raise AssertionError(
            "Skill manifest script inventory differs from scripts/*.py: "
            f"missing={sorted(actual_script_paths - manifest_script_paths)} "
            f"extra={sorted(manifest_script_paths - actual_script_paths)}"
        )

    manifest_paths = declared_paths(skill_manifest) | declared_paths(template_manifest)
    referenced_paths = {
        path
        for path in re.findall(r"`((?:assets|references|scripts)/[^`\s]+)`", skill_text)
        if "*" not in path
    }
    assert_paths_exist(manifest_paths, "Manifests")
    assert_paths_exist(referenced_paths, "SKILL.md")

    contracts = skill_manifest.get("contracts", {})
    evaluations = skill_manifest.get("evaluations", {})
    error_codes = skill_manifest.get("error_contract", {}).get("expected_cli_codes", [])
    if contracts.get("input_authority") != "references/rule-intake.md":
        raise AssertionError("Skill manifest must declare the input-authority contract")
    if evaluations.get("activation_safety_corpus") != "references/evals/activation-safety-regression.json":
        raise AssertionError("Skill manifest must declare the activation-safety corpus")
    if "E_INPUT_UNCONFIRMED" not in error_codes:
        raise AssertionError("Skill manifest must declare E_INPUT_UNCONFIRMED")
    return {
        "skill_version": EXPECTED_SKILL_VERSION,
        "template_bundle_version": EXPECTED_TEMPLATE_VERSION,
        "manifest_paths": len(manifest_paths),
        "script_inventory": len(actual_script_paths),
    }


def check_no_stale_contracts() -> dict[str, Any]:
    forbidden = {
        "V8.02": "stale version",
        "rule-layout-safety-v802": "stale layout rule",
        "/path/to/gen-tgo-ppt-skill": "placeholder Skill path",
        "all eight styles": "stale style count",
        "所有独立规则都必须有独立规则子智能体": "obsolete mandatory subagent rule",
        "为本次已加载的每个 rule 文件分配不同规则子智能体": "obsolete per-rule subagent rule",
        "不得让同一个子智能体代审多个 rule": "obsolete per-rule reviewer rule",
    }
    checked = 0
    for path in ROOT.rglob("*"):
        if not path.is_file() or path.suffix.lower() not in TEXT_EXTENSIONS:
            continue
        if ".git" in path.parts or "__pycache__" in path.parts or path.resolve() == Path(__file__).resolve():
            continue
        checked += 1
        text = path.read_text(encoding="utf-8")
        for needle, reason in forbidden.items():
            if needle in text:
                raise AssertionError(f"{reason}: {path.relative_to(ROOT)}")
    return {"text_files_checked": checked}


def check_activation_and_input_authority_contracts() -> dict[str, Any]:
    safety_path = ROOT / "references" / "evals" / "activation-safety-regression.json"
    safety = load_json(safety_path)
    if safety.get("schema_version") != ACTIVATION_SAFETY_SCHEMA:
        raise AssertionError("Activation-safety corpus schema mismatch")
    cases = safety.get("cases")
    if not isinstance(cases, list) or len(cases) != 5:
        raise AssertionError("Activation-safety corpus must contain exactly five cases")

    case_ids: set[str] = set()
    prompt_hashes: set[str] = set()
    for case in cases:
        if not isinstance(case, dict):
            raise AssertionError("Activation-safety cases must be objects")
        case_id = case.get("case_id")
        prompt = case.get("prompt")
        digest = case.get("prompt_sha256")
        if not isinstance(case_id, str) or not isinstance(prompt, str) or not isinstance(digest, str):
            raise AssertionError("Activation-safety case identity/prompt/hash is incomplete")
        if case_id in case_ids or digest in prompt_hashes:
            raise AssertionError("Activation-safety case IDs and prompt hashes must be unique")
        if hashlib.sha256(prompt.encode("utf-8")).hexdigest() != digest:
            raise AssertionError(f"Activation-safety prompt hash mismatch: {case_id}")
        case_ids.add(case_id)
        prompt_hashes.add(digest)

    by_id = {case["case_id"]: case for case in cases}
    activation_ids = {"activation-only-run-skill-001", "activation-only-name-skill-002"}
    for case_id in activation_ids:
        case = by_id.get(case_id, {})
        if case.get("expected_mode") is not None or case.get("expected_error") != "E_ROUTE_AMBIGUOUS":
            raise AssertionError(f"Activation-only case must stop at E_ROUTE_AMBIGUOUS: {case_id}")
        required_forbidden = {"scan_workspace", "read_workspace_content", "create_log", "generate_output"}
        if not required_forbidden.issubset(set(case.get("forbidden_actions", []))):
            raise AssertionError(f"Activation-only case lacks forbidden actions: {case_id}")

    unconfirmed = by_id.get("create-with-unconfirmed-workspace-003", {})
    if (
        unconfirmed.get("expected_mode") != "create"
        or unconfirmed.get("expected_error") != "E_INPUT_UNCONFIRMED"
        or "read_workspace_content" not in unconfirmed.get("forbidden_actions", [])
    ):
        raise AssertionError("Unconfirmed-workspace regression contract is incomplete")
    if by_id.get("create-with-explicit-workspace-files-004", {}).get("expected_source_authority") != "explicit_path":
        raise AssertionError("Explicit workspace paths must activate explicit_path authority")
    if by_id.get("create-with-current-turn-content-005", {}).get("expected_source_authority") != "current_turn":
        raise AssertionError("Current-turn content must activate current_turn authority")

    required_markers = {
        "SKILL.md": ["E_ROUTE_AMBIGUOUS", "本次已确认输入"],
        "references/rule-routing-and-modes.md": ["明确执行意图门", "不得读取这些文件的内容"],
        "references/rule-intake.md": ["本次输入归属门", "E_INPUT_UNCONFIRMED", "长期记忆"],
        "references/workflow.md": ["在扫描工作区之前执行明确意图门", "未确认的工作区文件不得读取正文"],
        "references/rule-guardrails.md": ["仅调用或点名 Skill 不得推断为 `create`", "不能自动成为本次输入"],
        "references/rule-failure-recovery.md": ["E_INPUT_UNCONFIRMED"],
    }
    for relative, markers in required_markers.items():
        text = (ROOT / relative).read_text(encoding="utf-8")
        missing = [marker for marker in markers if marker not in text]
        if missing:
            raise AssertionError(f"Activation/input contract markers missing from {relative}: {missing}")
    return {
        "safety_cases": len(cases),
        "activation_only_cases": len(activation_ids),
        "error_codes": ["E_ROUTE_AMBIGUOUS", "E_INPUT_UNCONFIRMED"],
    }


def check_template_assets() -> dict[str, Any]:
    manifest = load_json(TEMPLATE_MANIFEST_PATH)
    templates = manifest.get("templates")
    design_templates = manifest.get("design_resources", {}).get("templates")
    if not isinstance(templates, dict) or not isinstance(design_templates, dict):
        raise AssertionError("Template manifest must declare templates and design_resources.templates")
    if set(templates) != set(design_templates):
        raise AssertionError(
            "Selectable template keys differ from design routes: "
            f"templates={sorted(templates)} design_routes={sorted(design_templates)}"
        )

    evidence: dict[str, Any] = {}
    for key, record in sorted(templates.items()):
        if not isinstance(record, dict):
            raise AssertionError(f"Template record must be an object: {key}")
        asset = record.get("asset")
        if not isinstance(asset, str) or not asset.endswith(".pptx"):
            raise AssertionError(f"Template {key} must declare a PPTX asset")
        path = ROOT / asset
        if not path.is_file():
            raise AssertionError(f"Template asset is missing: {asset}")
        try:
            with ZipFile(path) as archive:
                names = archive.namelist()
        except Exception as exc:
            raise AssertionError(f"Template asset is not a readable PPTX package: {asset}: {exc}") from exc
        slide_count = sum(bool(re.fullmatch(r"ppt/slides/slide\d+\.xml", name)) for name in names)
        master_count = sum(bool(re.fullmatch(r"ppt/slideMasters/slideMaster\d+\.xml", name)) for name in names)
        if record.get("slides") != slide_count or record.get("masters") != master_count:
            raise AssertionError(
                f"Template metadata mismatch for {key}: "
                f"slides={slide_count}/{record.get('slides')} masters={master_count}/{record.get('masters')}"
            )
        expected_sha = record.get("asset_sha256")
        if expected_sha:
            actual_sha = hashlib.sha256(path.read_bytes()).hexdigest()
            if actual_sha != expected_sha:
                raise AssertionError(f"Template SHA-256 mismatch for {key}")
        evidence[key] = {"slides": slide_count, "masters": master_count, "asset": asset}

    loop = templates.get("loop-orange-white", {})
    menu_templates = manifest.get("menus", {}).get("template", [])
    branch = manifest.get("branch_styles", {}).get("D", {})
    if (
        loop.get("display_name") != "LOOP 大会"
        or loop.get("direct_select") is not True
        or "loop-orange-white" not in menu_templates
        or branch.get("asset") != loop.get("asset")
    ):
        raise AssertionError("LOOP conference template is not registered as a direct selectable asset")
    return {"template_count": len(evidence), "templates": evidence, "loop_direct_select": True}


def imported_roots(tree: ast.AST) -> set[str]:
    roots: set[str] = set()
    for node in ast.walk(tree):
        if isinstance(node, ast.Import):
            roots.update(alias.name.split(".", 1)[0] for alias in node.names)
        elif isinstance(node, ast.ImportFrom) and node.module and node.level == 0:
            roots.add(node.module.split(".", 1)[0])
    return roots


def check_ast_compile_and_dependency_graph() -> dict[str, Any]:
    manifest = load_json(SKILL_MANIFEST_PATH)
    runtime = manifest.get("runtime")
    if not isinstance(runtime, dict) or not isinstance(runtime.get("dependency_mapping"), dict):
        raise AssertionError("Skill manifest runtime.dependency_mapping is missing")
    if runtime["dependency_mapping"] != IMPORT_TO_DISTRIBUTION:
        raise AssertionError("Skill manifest import/distribution mapping mismatch")
    declared = {
        item["path"]: item for item in manifest.get("scripts", []) if isinstance(item, dict) and "path" in item
    }
    stdlib = set(sys.stdlib_module_names) | {"__future__"}
    graph: dict[str, list[str]] = {}
    for script in sorted(SCRIPTS.glob("*.py")):
        relative = script.relative_to(ROOT).as_posix()
        source = script.read_text(encoding="utf-8")
        tree = ast.parse(source, filename=relative)
        compile(tree, relative, "exec")
        third_party = sorted(imported_roots(tree) - stdlib)
        graph[relative] = third_party
        record = declared.get(relative)
        if record is None:
            raise AssertionError(f"Script missing from Skill manifest: {relative}")
        imports = record.get("imports")
        declared_imports = imports.get("third_party") if isinstance(imports, dict) else None
        if sorted(declared_imports or []) != third_party:
            raise AssertionError(
                f"Third-party import graph mismatch for {relative}: "
                f"manifest={declared_imports!r} AST={third_party!r}"
            )
        unknown = [name for name in third_party if name not in IMPORT_TO_DISTRIBUTION]
        if unknown:
            raise AssertionError(f"Undeclared import/distribution mapping for {relative}: {unknown}")
        expected_distributions = sorted({IMPORT_TO_DISTRIBUTION[name] for name in third_party})
        if sorted(record.get("distributions", [])) != expected_distributions:
            raise AssertionError(
                f"Distribution graph mismatch for {relative}: expected {expected_distributions}"
            )
    return {"scripts_compiled_in_memory": len(graph), "third_party_import_graph": graph}


def check_wrapper_contract(tree: ast.AST, source: str) -> None:
    if 'with_name("run_v11_skill_checks.py")' not in source:
        raise AssertionError("V1 wrapper must target run_v11_skill_checks.py exactly once")
    main = next(
        (node for node in tree.body if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef)) and node.name == "main"),
        None,
    )
    if main is None:
        raise AssertionError("V1 wrapper main() is missing")
    parse_lines: list[int] = []
    spawn_lines: list[int] = []
    for node in ast.walk(main):
        if isinstance(node, ast.Call) and isinstance(node.func, ast.Attribute):
            if node.func.attr == "parse_args":
                parse_lines.append(node.lineno)
            if node.func.attr == "run" and isinstance(node.func.value, ast.Name) and node.func.value.id == "subprocess":
                spawn_lines.append(node.lineno)
    if not parse_lines or not spawn_lines or min(parse_lines) >= min(spawn_lines):
        raise AssertionError("V1 wrapper must parse --help before spawning the V1.1 harness")
    if not any(
        isinstance(node, ast.Return)
        and isinstance(node.value, ast.Attribute)
        and node.value.attr == "returncode"
        for node in ast.walk(main)
    ):
        raise AssertionError("V1 wrapper must return the delegated process exit code")
    if ENTRYPOINT_ENV not in source:
        raise AssertionError("V1 wrapper must identify the compatibility entrypoint")


def check_cli_help_dependency_errors_and_wrapper_contract() -> dict[str, Any]:
    wrapper_source = COMPATIBILITY_WRAPPER.read_text(encoding="utf-8")
    check_wrapper_contract(ast.parse(wrapper_source, filename=str(COMPATIBILITY_WRAPPER)), wrapper_source)

    help_checks = 0
    for script in sorted(SCRIPTS.glob("*.py")):
        if script == COMPATIBILITY_WRAPPER:
            continue  # The canonical harness must never invoke its compatibility wrapper.
        for flags in (["-B"], ["-I", "-S", "-B"]):
            result = run([sys.executable, *flags, str(script), "--help"])
            if "usage:" not in result.stdout.lower():
                raise AssertionError(f"--help output lacks usage for {script.name} with flags {flags}")
            help_checks += 1

    dependency_cases = {
        "check_pptx_layout.py": ["python-pptx"],
        "inspect_pptx_style.py": ["lxml", "Pillow", "python-pptx"],
    }
    missing_input = ROOT / ".gen-tgo-ppt-harness-missing-input.pptx"
    if missing_input.exists():
        raise AssertionError(f"Reserved missing-input probe unexpectedly exists: {missing_input}")
    for script_name, distributions in dependency_cases.items():
        result = run(
            [sys.executable, "-I", "-S", "-B", str(SCRIPTS / script_name), str(missing_input)],
            expected=2,
        )
        combined = result.stdout + result.stderr
        if "E_DEPENDENCY_MISSING" not in result.stderr or "Traceback" in combined:
            raise AssertionError(f"{script_name} dependency error contract is unstable")
        for distribution in distributions:
            if distribution not in result.stderr:
                raise AssertionError(f"{script_name} dependency error omits {distribution}")
    return {
        "help_invocations": help_checks,
        "isolated_help": True,
        "compatibility_wrapper": "AST-verified; intentionally not invoked",
        "dependency_error_code": "E_DEPENDENCY_MISSING",
    }


def check_generation_log_dry_run_and_actual() -> dict[str, Any]:
    script = SCRIPTS / "create_generation_log.py"
    with temporary_directory() as temporary:
        cwd = Path(temporary)
        rejected = run(
            [sys.executable, "-B", str(script), "--mode", "create", "--title", "Harness 未确认输入"],
            cwd=cwd,
            expected=1,
        )
        if "E_INPUT_UNCONFIRMED" not in rejected.stderr or list(cwd.iterdir()):
            raise AssertionError("Generation log must reject unconfirmed create input without writing")
        dry = run(
            [
                sys.executable,
                "-B",
                str(script),
                "--dry-run",
                "--mode",
                "create",
                "--source-authority",
                "current_turn",
                "--title",
                "Harness 预览",
            ],
            cwd=cwd,
        )
        if list(cwd.iterdir()):
            raise AssertionError("Generation-log --dry-run wrote files")
        for required in ("Skill 版本：V1.1", "排版安全版本：V1", "模板包版本：V1.1", "任务模式：create", "输入归属：current_turn", "## 失败、重试与降级"):
            if required not in dry.stdout:
                raise AssertionError(f"Generation-log dry-run is missing {required!r}")

        actual = run(
            [
                sys.executable,
                "-B",
                str(script),
                "--mode",
                "convert",
                "--source-authority",
                "explicit_path",
                "--title",
                "Harness 实写",
            ],
            cwd=cwd,
        )
        output_path = Path(actual.stdout.strip())
        if not output_path.is_absolute():
            output_path = cwd / output_path
        if not output_path.is_file() or output_path.parent.resolve() != cwd.resolve():
            raise AssertionError("Generation-log actual mode did not create the announced file")
        files = list(cwd.glob("gen-tgo-ppt-生成日志-*.md"))
        if len(files) != 1 or files[0].resolve() != output_path.resolve():
            raise AssertionError("Generation-log actual mode must create exactly one log")
        text = output_path.read_text(encoding="utf-8")
        for required in ("Skill 版本：V1.1", "排版安全版本：V1", "模板包版本：V1.1", "任务模式：convert", "输入归属：explicit_path", "## 遗留问题"):
            if required not in text:
                raise AssertionError(f"Generated log is missing {required!r}")
    return {"unconfirmed_rejected": 1, "dry_run_writes": 0, "actual_logs_created": 1}


def check_html_fixtures() -> dict[str, Any]:
    good = """<!doctype html>
<html><head><style>
.slide { aspect-ratio: 16 / 9; width: 1600px; height: 900px; }
@media print { .slide { page-break-after: always; } }
</style></head><body>
<section class="slide"><h1>标题</h1></section>
<section class="slide"><h1>嘉宾介绍</h1></section>
<section class="slide"><h1>目录</h1><p>一、背景</p></section>
<section class="slide"><h1>感谢聆听</h1></section>
</body></html>"""
    bad = """<!doctype html><html><body><section class="slide"><p>""" + ("长文本" * 700) + """</p></section></body></html>"""
    with temporary_directory() as temporary:
        root = Path(temporary)
        good_path = root / "good.html"
        bad_path = root / "bad.html"
        good_path.write_text(good, encoding="utf-8")
        bad_path.write_text(bad, encoding="utf-8")
        good_result = run([sys.executable, "-B", str(SCRIPTS / "check_html_layout.py"), str(good_path), "--json"])
        good_report = parse_json_stdout(good_result, "good HTML fixture")[0]
        if good_report.get("status") != "PASS" or good_report.get("slide_count") != 4:
            raise AssertionError(f"Good HTML fixture failed: {good_report}")
        bad_result = run(
            [sys.executable, "-B", str(SCRIPTS / "check_html_layout.py"), str(bad_path), "--json"],
            expected=1,
        )
        bad_report = parse_json_stdout(bad_result, "bad HTML fixture")[0]
        codes = {item.get("code") for item in bad_report.get("issues", [])}
        required_codes = {"CANVAS_16_9_MISSING", "MANDATORY_GUEST_PAGE_MISSING", "TEXT_DENSITY_HIGH"}
        if bad_report.get("status") != "FAIL" or not required_codes.issubset(codes):
            raise AssertionError(f"Bad HTML fixture did not expose stable failures: {bad_report}")
    return {"good_status": "PASS", "bad_status": "FAIL", "bad_issue_codes": sorted(required_codes)}


def pptx_dependencies_available() -> bool:
    return all(importlib.util.find_spec(name) is not None for name in ("pptx", "lxml", "PIL"))


def check_pptx_fixtures() -> dict[str, Any]:
    if not pptx_dependencies_available():
        return {"fixture_status": "SKIPPED", "reason": "optional PPTX distributions unavailable"}
    generator = r'''from pathlib import Path
import sys
from pptx import Presentation
from pptx.util import Inches, Pt

def build(path, texts):
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    layout = deck.slide_layouts[6]
    for text in texts:
        slide = deck.slides.add_slide(layout)
        if text:
            box = slide.shapes.add_textbox(Inches(2.0), Inches(2.0), Inches(8.0), Inches(1.0))
            box.text_frame.word_wrap = False
            paragraph = box.text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = text
            run.font.name = "Arial"
            run.font.size = Pt(24)
    deck.save(path)

build(Path(sys.argv[1]), ["V1.1 Harness", "嘉宾介绍", "感谢聆听"])
build(Path(sys.argv[2]), ["Only a cover"])
'''
    checker_flags = [
        "--json",
        "--no-check-key-page-solid-panel",
        "--no-check-agenda-page",
        "--disable-cover-safe-zone",
        "--allow-theme-fonts-on-key-pages",
        "--margin-in", "0",
        "--footer-in", "0",
        "--logo-width-in", "0",
        "--logo-height-in", "0",
    ]
    with temporary_directory() as temporary:
        root = Path(temporary)
        good_path = root / "good.pptx"
        bad_path = root / "bad.pptx"
        run([sys.executable, "-B", "-c", generator, str(good_path), str(bad_path)])
        good_result = run(
            [sys.executable, "-B", str(SCRIPTS / "check_pptx_layout.py"), str(good_path), *checker_flags]
        )
        good_report = parse_json_stdout(good_result, "good PPTX fixture")[0]
        if good_report.get("status") != "PASS" or good_report.get("slide_count") != 3:
            raise AssertionError(f"Good PPTX fixture failed: {good_report}")
        bad_result = run(
            [sys.executable, "-B", str(SCRIPTS / "check_pptx_layout.py"), str(bad_path), *checker_flags],
            expected=1,
        )
        bad_report = parse_json_stdout(bad_result, "bad PPTX fixture")[0]
        codes = {item.get("code") for item in bad_report.get("issues", [])}
        required_codes = {"MANDATORY_GUEST_PAGE_MISSING", "MANDATORY_CLOSING_MISSING"}
        if bad_report.get("status") != "FAIL" or not required_codes.issubset(codes):
            raise AssertionError(f"Bad PPTX fixture did not expose stable failures: {bad_report}")
        inspect_result = run(
            [sys.executable, "-B", str(SCRIPTS / "inspect_pptx_style.py"), str(good_path)]
        )
        inspect_report = parse_json_stdout(inspect_result, "PPTX style fixture")
        if len(inspect_report) != 1 or inspect_report[0].get("slide_count") != 3:
            raise AssertionError("PPTX style inspector did not parse the synthetic good deck")
    return {"fixture_status": "EXECUTED", "good_status": "PASS", "bad_issue_codes": sorted(required_codes)}


def check_helper_contracts() -> dict[str, Any]:
    with temporary_directory() as temporary:
        root = Path(temporary)
        low = root / "low.md"
        high = root / "high.md"
        low.write_text("# 简介\n\n一句简短说明。\n", encoding="utf-8")
        high.write_text("# 高密度\n\n" + "高密度文本" * 300, encoding="utf-8")
        budget = run(
            [sys.executable, "-B", str(SCRIPTS / "content_budget.py"), str(low), str(high), "--format", "json"]
        )
        budget_report = parse_json_stdout(budget, "content budget fixtures")
        risks = {Path(item["file"]).name: item["sections"][0]["risk"] for item in budget_report}
        if risks != {"low.md": "low", "high.md": "high"}:
            raise AssertionError(f"Content-budget risk contract changed: {risks}")

        mixed_path = root / "layout.json"
        mixed_path.write_text(
            json.dumps(
                [
                    {"file": "good.pptx", "status": "PASS", "slide_count": 3, "issues": []},
                    {
                        "file": "bad.pptx",
                        "status": "FAIL",
                        "slide_count": 1,
                        "issues": [
                            {"level": "FAIL", "slide": 1, "code": "MANDATORY_CLOSING_MISSING", "message": "missing"},
                            {"level": "WARN", "slide": 1, "code": "FONT_SMALL", "message": "small"},
                        ],
                    },
                ],
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        summary = run(
            [sys.executable, "-B", str(SCRIPTS / "summarize_layout_report.py"), str(mixed_path), "--format", "json"]
        )
        summary_report = parse_json_stdout(summary, "mixed layout summary")
        if summary_report.get("deck_count") != 2 or summary_report.get("status_counts") != {"PASS": 1, "FAIL": 1}:
            raise AssertionError("Layout-summary aggregate contract changed")
        if summary_report.get("issue_code_counts", {}).get("MANDATORY_CLOSING_MISSING") != 1:
            raise AssertionError("Layout-summary issue counts are incorrect")

        scan_root = root / "workspace"
        scan_root.mkdir()
        (scan_root / "Design.md").write_text("# Design\n", encoding="utf-8")
        (scan_root / "Content.md").write_text("# Content\n", encoding="utf-8")
        (scan_root / "gen-tgo-ppt-生成日志-20260716-010101.md").write_text("# Log\n", encoding="utf-8")
        (scan_root / "gen-tgo-ppt-SSOT-20260716-010101.md").write_text("# SSOT\n", encoding="utf-8")
        (scan_root / "source.md").write_text("# Source\n", encoding="utf-8")
        (scan_root / "tgo-template-contact-sheet.png").write_bytes(b"fixture")
        (scan_root / "brand-logo.png").write_bytes(b"fixture")
        scan = run(
            [sys.executable, "-B", str(SCRIPTS / "scan_task_context.py"), str(scan_root), "--recent-limit", "20"]
        )
        scan_report = parse_json_stdout(scan, "task-context fixture")
        source_names = {item["name"] for item in scan_report["categories"]["source_candidate"]["recent"]}
        logo_names = {item["name"] for item in scan_report["categories"]["logo_candidate"]["recent"]}
        if source_names != {"source.md"} or logo_names != {"brand-logo.png"}:
            raise AssertionError(f"Task-context noise filtering changed: sources={source_names}, logos={logo_names}")
        authority = scan_report.get("input_authority", {})
        if (
            authority.get("status") != "unknown_to_scanner"
            or authority.get("workspace_files_are_candidates_only") is not True
            or authority.get("requires_external_confirmation") is not True
            or scan_report.get("status_flags", {}).get("has_unconfirmed_workspace_candidates") is not True
        ):
            raise AssertionError("Task-context scanner must label workspace files as unconfirmed candidates")
    return {
        "content_budget": risks,
        "layout_decks": 2,
        "scan_sources": sorted(source_names),
        "workspace_authority": "unknown_to_scanner",
    }


def assert_evaluator_result(
    result: subprocess.CompletedProcess[str],
    *,
    status: str,
    exit_code: int,
    label: str,
) -> dict[str, Any]:
    if result.returncode != exit_code:
        raise AssertionError(f"{label} returned {result.returncode}, expected {exit_code}: {result.stdout}")
    if "Traceback" in result.stdout + result.stderr:
        raise AssertionError(f"{label} emitted a raw traceback")
    report = parse_json_stdout(result, label)
    if report.get("status") != status or report.get("exit_code") != exit_code:
        raise AssertionError(f"{label} status contract mismatch: {report}")
    expected_classification = {
        "PASS": "fixture_valid",
        "FAIL": "fixture_valid",
        "INVALID": "invalid",
        "INSUFFICIENT": "insufficient",
    }[status]
    if report.get("evidence_classification") != expected_classification:
        raise AssertionError(
            f"{label} evidence classification mismatch: "
            f"expected {expected_classification!r}, got {report.get('evidence_classification')!r}"
        )
    return report


def run_evaluator(corpus: Path, observations: Path, *extra: str) -> subprocess.CompletedProcess[str]:
    return run(
        [
            sys.executable,
            "-B",
            str(SCRIPTS / "evaluate_v11_contract.py"),
            str(corpus),
            str(observations),
            *extra,
        ],
        expected=None,
    )


def check_evaluator_fixtures_and_negative_paths() -> dict[str, Any]:
    eval_root = ROOT / "references" / "evals"
    paths = {
        split: (
            eval_root / f"routing-{split}.json",
            eval_root / f"evaluator-smoke-observations-{split}.json",
        )
        for split in ("regression", "holdout")
    }
    all_hashes: dict[str, set[str]] = {}
    for split, (corpus_path, observations_path) in paths.items():
        corpus = load_json(corpus_path)
        observations = load_json(observations_path)
        if corpus.get("schema_version") != EVAL_SCHEMA or corpus.get("split") != split:
            raise AssertionError(f"{split} corpus schema/split mismatch")
        cases = corpus.get("cases")
        if not isinstance(cases, list) or len(cases) != 6:
            raise AssertionError(f"{split} corpus must contain exactly six cases")
        if {case.get("case_type") for case in cases} != CASE_TYPES:
            raise AssertionError(f"{split} corpus does not cover all six case types")
        hashes = {case.get("prompt_sha256") for case in cases}
        if len(hashes) != 6:
            raise AssertionError(f"{split} corpus prompt hashes are not unique")
        for case in cases:
            expected_hash = hashlib.sha256(case["prompt"].encode("utf-8")).hexdigest()
            if case.get("prompt_sha256") != expected_hash:
                raise AssertionError(f"Prompt hash mismatch for {case.get('case_id')}")
        all_hashes[split] = hashes
        if observations.get("evidence_level") != "structural_fixture":
            raise AssertionError("Fixture observations must be labeled structural_fixture")
        serialized_observations = json.dumps(observations, ensure_ascii=False)
        if re.search(r'"(?:expected[^\"]*|gold[^\"]*)"\s*:', serialized_observations, flags=re.I):
            raise AssertionError("Fixture observations leak expected/gold fields")
        report = assert_evaluator_result(
            run_evaluator(corpus_path, observations_path),
            status="PASS",
            exit_code=0,
            label=f"{split} evaluator fixture",
        )
        if report.get("evidence_level") != "structural_fixture":
            raise AssertionError("Evaluator fixture report overstates behavior evidence")
        if report.get("confusion_matrix") != {"tp": 5, "fp": 0, "tn": 1, "fn": 0}:
            raise AssertionError(f"Evaluator active/no-trigger confusion contract changed: {report}")
        for metric in ("precision", "recall", "mode_accuracy", "specificity", "handoff_accuracy"):
            if report.get("metrics", {}).get(metric, {}).get("value") != 1.0:
                raise AssertionError(f"Evaluator fixture metric {metric} must be 1.0")
        if report.get("sample_scope", {}).get("fixture_results_are_agent_behavior_evidence") is not False:
            raise AssertionError("Evaluator fixture report must deny Agent-behavior evidence")
    if all_hashes["regression"] & all_hashes["holdout"]:
        raise AssertionError("Regression and holdout corpora share prompt hashes")

    regression_corpus, regression_observations = paths["regression"]
    base_observations = load_json(regression_observations)
    with temporary_directory() as temporary:
        root = Path(temporary)

        missing = copy.deepcopy(base_observations)
        missing["observations"] = missing["observations"][:-1]
        missing_path = root / "missing-observation.json"
        missing_path.write_text(json.dumps(missing, ensure_ascii=False), encoding="utf-8")
        assert_evaluator_result(
            run_evaluator(regression_corpus, missing_path),
            status="INSUFFICIENT",
            exit_code=3,
            label="missing-observation negative path",
        )

        duplicate = copy.deepcopy(base_observations)
        duplicate["observations"].append(copy.deepcopy(duplicate["observations"][0]))
        duplicate_path = root / "duplicate-observation.json"
        duplicate_path.write_text(json.dumps(duplicate, ensure_ascii=False), encoding="utf-8")
        assert_evaluator_result(
            run_evaluator(regression_corpus, duplicate_path),
            status="INVALID",
            exit_code=2,
            label="duplicate-observation negative path",
        )

        invalid = copy.deepcopy(base_observations)
        invalid["observations"][0]["predicted_route"] = "invalid-route"
        invalid_path = root / "invalid-label.json"
        invalid_path.write_text(json.dumps(invalid, ensure_ascii=False), encoding="utf-8")
        assert_evaluator_result(
            run_evaluator(regression_corpus, invalid_path),
            status="INVALID",
            exit_code=2,
            label="invalid-label negative path",
        )

        threshold = copy.deepcopy(base_observations)
        threshold["observations"][0]["predicted_route"] = "no_trigger"
        threshold["observations"][0]["predicted_mode"] = None
        threshold_path = root / "threshold-fail.json"
        threshold_path.write_text(json.dumps(threshold, ensure_ascii=False), encoding="utf-8")
        assert_evaluator_result(
            run_evaluator(regression_corpus, threshold_path),
            status="FAIL",
            exit_code=1,
            label="threshold-fail negative path",
        )

        assert_evaluator_result(
            run_evaluator(
                regression_corpus,
                regression_observations,
                "--case-id",
                "regression-no-trigger-001",
            ),
            status="INSUFFICIENT",
            exit_code=3,
            label="zero-positive negative path",
        )
    return {
        "fixture_splits": ["regression", "holdout"],
        "cases_per_split": 6,
        "evidence_level": "structural_fixture",
        "negative_paths": ["missing", "duplicate", "invalid_label", "threshold_fail", "zero_positive"],
    }


def check_packaging() -> dict[str, Any]:
    nested_git = [path for path in ROOT.rglob(".git") if path.is_dir() and path.parent != ROOT]
    if nested_git:
        raise AssertionError(f"Nested .git directories must not be packaged: {nested_git}")
    cache_dirs = [path for path in ROOT.rglob("__pycache__") if path.is_dir()]
    compiled_files = [path for path in ROOT.rglob("*.py[co]") if path.is_file()]
    if cache_dirs or compiled_files:
        raise AssertionError(
            f"Python cache artifacts must not be packaged: caches={cache_dirs}, compiled={compiled_files}"
        )
    broken_links = [path for path in ROOT.rglob("*") if path.is_symlink() and not path.exists()]
    if broken_links:
        raise AssertionError(f"Broken symlinks must not be packaged: {broken_links}")
    readme = (ROOT / "README.md").read_text(encoding="utf-8")
    required = ["如何增加新模板", "references/skill-manifest.json", "run_v11_skill_checks.py", "run_v1_skill_checks.py"]
    missing = [item for item in required if item not in readme]
    if missing:
        raise AssertionError(f"README.md packaging/maintenance instructions are incomplete: {missing}")
    return {"nested_git_dirs": 0, "python_cache_artifacts": 0, "broken_symlinks": 0}


CHECKS: list[tuple[str, str, Callable[[], dict[str, Any]]]] = [
    ("frontmatter_and_protected_section", "structural", check_frontmatter_and_protected_section),
    ("version_manifests_agents_and_paths", "structural", check_version_manifests_agents_and_paths),
    ("template_assets", "structural", check_template_assets),
    ("no_stale_contracts", "structural", check_no_stale_contracts),
    ("activation_and_input_authority_contracts", "structural_fixture", check_activation_and_input_authority_contracts),
    ("ast_compile_and_dependency_graph", "structural", check_ast_compile_and_dependency_graph),
    ("cli_help_dependency_errors_and_wrapper_contract", "deterministic_contract", check_cli_help_dependency_errors_and_wrapper_contract),
    ("generation_log_dry_run_and_actual", "deterministic_fixture", check_generation_log_dry_run_and_actual),
    ("html_fixtures", "deterministic_fixture", check_html_fixtures),
    ("pptx_fixtures", "deterministic_fixture", check_pptx_fixtures),
    ("helper_contracts", "deterministic_fixture", check_helper_contracts),
    ("evaluator_fixtures_and_negative_paths", "structural_fixture", check_evaluator_fixtures_and_negative_paths),
    ("packaging", "structural", check_packaging),
]


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Run V1.1 structural, deterministic CLI, helper, and evaluator-fixture checks."
    )
    parser.add_argument("--compact", action="store_true", help="Emit compact JSON instead of indented JSON.")
    parser.add_argument(
        "--temp-root",
        type=Path,
        default=None,
        help=(
            "Existing writable directory for isolated fixture workspaces. "
            f"Defaults to ${TEMP_ROOT_ENV}, then the system temporary directory."
        ),
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    global TEMP_ROOT
    args = build_parser().parse_args(argv)
    TEMP_ROOT, temp_workspace = resolve_temp_root(args.temp_root)
    results: list[dict[str, Any]] = []
    for name, evidence_level, check in CHECKS:
        if name in TEMP_REQUIRED_CHECKS and TEMP_ROOT is None:
            results.append(
                {
                    "check": name,
                    "status": "INSUFFICIENT",
                    "evidence_level": evidence_level,
                    "error_code": "E_TEMP_ROOT_UNAVAILABLE",
                    "error": "deterministic fixture gate was not executed because no writable temp root is available",
                }
            )
            continue
        try:
            evidence = check()
        except Exception as exc:
            results.append(
                {
                    "check": name,
                    "status": "FAIL",
                    "evidence_level": evidence_level,
                    "error": f"{type(exc).__name__}: {exc}",
                }
            )
        else:
            results.append(
                {
                    "check": name,
                    "status": "PASS",
                    "evidence_level": evidence_level,
                    "evidence": evidence,
                }
            )

    if any(item["status"] == "FAIL" for item in results):
        status = "FAIL"
    elif any(item["status"] == "INSUFFICIENT" for item in results):
        status = "INSUFFICIENT"
    else:
        status = "PASS"
    entrypoint = os.environ.get(ENTRYPOINT_ENV, Path(__file__).name)
    if entrypoint not in {"run_v11_skill_checks.py", "run_v1_skill_checks.py"}:
        entrypoint = Path(__file__).name
    payload = {
        "status": status,
        "entrypoint": entrypoint,
        "harness_version": EXPECTED_SKILL_VERSION,
        "evidence_level": "structural_and_deterministic_fixture",
        "temp_workspace": temp_workspace,
        "checks": results,
    }
    print(json.dumps(payload, ensure_ascii=False, indent=None if args.compact else 2, sort_keys=args.compact))
    return {"PASS": 0, "FAIL": 1, "INSUFFICIENT": 3}[status]


if __name__ == "__main__":
    raise SystemExit(main())

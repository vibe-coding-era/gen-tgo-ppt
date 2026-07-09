#!/usr/bin/env python3
"""Static PPTX layout checker for gen-tgo-ppt V1."""

from __future__ import annotations

import argparse
import json
import math
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_IN = 914400
EMU_PER_PT = 12700
BASE_SLIDE_WIDTH_IN = 26.667
BASE_SLIDE_HEIGHT_IN = 15.0
PROJECT_INFO_LABELS = {
    "大会主题",
    "会议主题",
    "活动主题",
    "项目名称",
    "项目主题",
    "时间",
    "地点",
    "时间地点",
    "参会人群",
    "目标人群",
    "合作身份",
    "赞助身份",
    "负责人",
    "主办方",
    "承办方",
}
AGENDA_MARKERS = {"目录", "议程", "大纲", "contents", "agenda"}
DEFAULT_AGENDA_FORBIDDEN_PATTERNS = (
    "今天只讲",
    "只讲一个问题",
)


@dataclass(frozen=True)
class Rect:
    left: int
    top: int
    width: int
    height: int

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height

    @property
    def area(self) -> int:
        return max(0, self.width) * max(0, self.height)


def inch_to_emu(value: float) -> int:
    return int(value * EMU_PER_IN)


def emu_to_in(value: int) -> float:
    return round(value / EMU_PER_IN, 3)


def emu_to_pt(value: int) -> float:
    return value / EMU_PER_PT


def resolve_layout_scale(slide_w: int, slide_h: int, args: argparse.Namespace) -> float:
    if args.layout_scale is not None:
        return args.layout_scale
    if args.no_auto_scale:
        return 1.0

    width_scale = slide_w / EMU_PER_IN / BASE_SLIDE_WIDTH_IN
    height_scale = slide_h / EMU_PER_IN / BASE_SLIDE_HEIGHT_IN
    scale = min(width_scale, height_scale)
    if abs(scale - 1.0) < 0.03:
        return 1.0
    if 0.25 <= scale <= 1.05:
        return scale
    return 1.0


def scaled(value: float, layout_scale: float) -> float:
    return value * layout_scale


def shape_rect(shape) -> Rect:
    return Rect(
        int(getattr(shape, "left", 0) or 0),
        int(getattr(shape, "top", 0) or 0),
        int(getattr(shape, "width", 0) or 0),
        int(getattr(shape, "height", 0) or 0),
    )


def intersects(a: Rect, b: Rect) -> bool:
    return a.left < b.right and a.right > b.left and a.top < b.bottom and a.bottom > b.top


def intersection_area(a: Rect, b: Rect) -> int:
    if not intersects(a, b):
        return 0
    width = min(a.right, b.right) - max(a.left, b.left)
    height = min(a.bottom, b.bottom) - max(a.top, b.top)
    return max(0, width) * max(0, height)


def walk_shapes(shapes):
    for shape in shapes:
        yield shape
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk_shapes(shape.shapes)
        except Exception:
            pass


def shape_text(shape) -> str:
    try:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            return ""
        return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
    except Exception:
        return ""


def normalize_text(text: str) -> str:
    return "".join(text.split()).casefold()


def text_preview(text: str, limit: int = 28) -> str:
    compact = " ".join(text.split())
    if len(compact) <= limit:
        return compact
    return compact[: limit - 1] + "..."


def compact_text(text: str) -> str:
    return " ".join(text.split())


def is_agenda_marker_text(text: str) -> bool:
    normalized = normalize_text(text)
    return any(normalize_text(marker) in normalized for marker in AGENDA_MARKERS)


def slide_layout_name(slide) -> str:
    try:
        return str(slide.slide_layout.name or "")
    except Exception:
        return ""


def is_agenda_slide(slide, texts: list[str]) -> bool:
    return is_agenda_marker_text(slide_layout_name(slide)) or any(is_agenda_marker_text(text) for text in texts)


def parse_rgb(value: str) -> tuple[int, int, int]:
    compact = value.strip().lstrip("#")
    if len(compact) != 6:
        raise ValueError(f"Expected 6-digit RGB hex, got {value!r}")
    return (int(compact[0:2], 16), int(compact[2:4], 16), int(compact[4:6], 16))


def shape_fill_rgb(shape) -> Optional[tuple[int, int, int]]:
    try:
        rgb = shape.fill.fore_color.rgb
    except Exception:
        return None
    if rgb is None:
        return None
    try:
        return parse_rgb(str(rgb))
    except ValueError:
        return None


def rgb_near(candidate: tuple[int, int, int], target: tuple[int, int, int], tolerance: int) -> bool:
    return max(abs(candidate[index] - target[index]) for index in range(3)) <= tolerance


def shape_label(shape, text: str) -> str:
    name = (getattr(shape, "name", "") or "unnamed").strip()
    preview = text_preview(text)
    if preview:
        return f"{name!r} text={preview!r}"
    return repr(name)


def font_sizes(shape, default_size: float) -> list[float]:
    sizes: list[float] = []
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is not None:
                    sizes.append(float(run.font.size.pt))
    except Exception:
        pass
    return sizes or [default_size]


def explicit_font_names(shape) -> list[str]:
    names: list[str] = []
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.name:
                    names.append(str(run.font.name))
    except Exception:
        pass
    return names


def vertical_text_padding_pt(shape, layout_scale: float) -> float:
    fallback = 10.0 * layout_scale
    try:
        text_frame = shape.text_frame
        margin_top = int(text_frame.margin_top or 0)
        margin_bottom = int(text_frame.margin_bottom or 0)
        explicit = emu_to_pt(margin_top + margin_bottom)
        return max(fallback, explicit)
    except Exception:
        return fallback


def weighted_text_len(text: str) -> float:
    total = 0.0
    for char in text:
        if char.isspace():
            total += 0.35
        elif ord(char) > 127:
            total += 1.0
        else:
            total += 0.55
    return total


def estimate_lines(text: str, width_emu: int, font_size_pt: float) -> int:
    available_width_pt = max(1.0, emu_to_pt(width_emu) - 16.0)
    units_per_line = max(1.0, available_width_pt / max(font_size_pt * 0.88, 1.0))
    lines = 0
    for paragraph in text.splitlines() or [text]:
        weighted = weighted_text_len(paragraph)
        if weighted <= 0:
            lines += 1
        else:
            lines += max(1, math.ceil(weighted / units_per_line))
    return lines


def likely_single_line_info(text: str, max_weighted_len: float) -> bool:
    if "\n" in text:
        return False
    compact = compact_text(text)
    if not compact:
        return False
    if any(mark in compact for mark in ("。", "；", ";", "！", "!", "？", "?")):
        return False
    if weighted_text_len(compact) > max_weighted_len:
        return False
    return any(token in compact for token in ("：", ":", "｜", "|"))


def project_info_label(text: str) -> str:
    compact = compact_text(text)
    for separator in ("：", ":"):
        if separator in compact:
            return compact.split(separator, 1)[0].strip()
    return ""


def is_project_info_item(text: str) -> bool:
    label = project_info_label(text)
    return bool(label and label in PROJECT_INFO_LABELS)


def text_word_wrap(shape):
    try:
        return shape.text_frame.word_wrap
    except Exception:
        return None


def level_rank(level: str) -> int:
    return {"PASS": 0, "WARN": 1, "FAIL": 2}[level]


def add_issue(issues: list[dict], level: str, slide: int, code: str, message: str) -> None:
    issues.append({"level": level, "slide": slide, "code": code, "message": message})


def check_mandatory_slide_content(
    issues: list[dict],
    slide_index: int,
    slide_count: int,
    texts: list[str],
    args: argparse.Namespace,
) -> None:
    if args.skip_mandatory_page_check:
        return

    normalized_texts = [normalize_text(text) for text in texts if normalize_text(text)]
    if slide_index == 2:
        guest_title = normalize_text(args.guest_page_title)
        if guest_title not in normalized_texts:
            add_issue(
                issues,
                "FAIL",
                slide_index,
                "MANDATORY_GUEST_TITLE_MISSING",
                f"Page 2 must contain the mandatory '{args.guest_page_title}' title.",
            )
        extra_texts = [
            text_preview(text, limit=36)
            for text in texts
            if normalize_text(text) and normalize_text(text) != guest_title
        ]
        if extra_texts and not args.allow_guest_page_content:
            add_issue(
                issues,
                "FAIL",
                slide_index,
                "MANDATORY_GUEST_PAGE_NOT_BLANK",
                "Page 2 must be a blank guest-introduction page; extra text found: "
                + ", ".join(repr(text) for text in extra_texts[:4])
                + ("." if len(extra_texts) <= 4 else ", ..."),
            )

    if slide_index == slide_count:
        closing_text = normalize_text(args.closing_text)
        slide_text = normalize_text("".join(texts))
        if closing_text not in slide_text and not args.allow_custom_closing:
            add_issue(
                issues,
                "FAIL",
                slide_index,
                "MANDATORY_CLOSING_MISSING",
                f"Final slide must contain the mandatory closing copy '{args.closing_text}'.",
            )


def check_cover_text_safe_zone(
    issues: list[dict],
    slide_index: int,
    shape,
    text: str,
    rect: Rect,
    slide_w: int,
    args: argparse.Namespace,
) -> None:
    if args.disable_cover_safe_zone or slide_index != 1:
        return
    safe_right = int(slide_w * args.cover_text_right_ratio)
    rect_center = rect.left + rect.width / 2
    is_left_column_text = rect_center < slide_w * args.cover_left_column_center_ratio
    if is_left_column_text and rect.left < safe_right and rect.right > safe_right:
        add_issue(
            issues,
            "FAIL",
            slide_index,
            "COVER_MAIN_VISUAL_KEEP_OUT",
            f"{shape_label(shape, text)} crosses the cover text-safe boundary at "
            f"x={emu_to_in(safe_right)} in; shorten, wrap inside the left column, "
            "or switch cover layout.",
        )


def looks_like_key_page_solid_panel(shape, slide_w: int, slide_h: int, target_rgb: tuple[int, int, int], args: argparse.Namespace) -> bool:
    rect = shape_rect(shape)
    if rect.width <= 0 or rect.height <= 0:
        return False

    rgb = shape_fill_rgb(shape)
    if rgb is None or not rgb_near(rgb, target_rgb, args.key_panel_color_tolerance):
        return False

    width_ratio = rect.width / max(1, slide_w)
    height_ratio = rect.height / max(1, slide_h)
    area_ratio = rect.area / max(1, slide_w * slide_h)
    if width_ratio < args.key_panel_min_width_ratio or width_ratio > args.key_panel_max_width_ratio:
        return False
    if height_ratio < args.key_panel_min_height_ratio or height_ratio > args.key_panel_max_height_ratio:
        return False
    if area_ratio < args.key_panel_min_area_ratio:
        return False

    center_x = (rect.left + rect.width / 2) / max(1, slide_w)
    center_y = (rect.top + rect.height / 2) / max(1, slide_h)
    if abs(center_x - 0.5) > args.key_panel_center_x_tolerance:
        return False
    if abs(center_y - 0.5) > args.key_panel_center_y_tolerance:
        return False
    return True


def has_key_page_solid_panel(slide, slide_w: int, slide_h: int, target_rgb: tuple[int, int, int], args: argparse.Namespace) -> bool:
    return any(
        looks_like_key_page_solid_panel(shape, slide_w, slide_h, target_rgb, args)
        for shape in walk_shapes(slide.shapes)
    )


def check_key_page_solid_panel(
    issues: list[dict],
    slide_index: int,
    slide_count: int,
    slide,
    slide_w: int,
    slide_h: int,
    args: argparse.Namespace,
) -> None:
    if not args.check_key_page_solid_panel:
        return
    if slide_index not in (1, slide_count):
        return
    if slide_index == 1 and args.allow_native_cover_layout:
        return
    if slide_index == slide_count and args.allow_native_closing_layout:
        return

    target_rgb = parse_rgb(args.key_panel_color)
    if has_key_page_solid_panel(slide, slide_w, slide_h, target_rgb, args):
        return

    role = "Title page" if slide_index == 1 else "Final closing page"
    exception_flag = "--allow-native-cover-layout" if slide_index == 1 else "--allow-native-closing-layout"
    add_issue(
        issues,
        "FAIL",
        slide_index,
        "KEY_PAGE_SOLID_PANEL_MISSING",
        f"{role} must use a centered opaque solid content panel near #{args.key_panel_color}; "
        f"use {exception_flag} only when the user explicitly approved the native layout exception.",
    )


def check_agenda_page_order(
    issues: list[dict],
    slide_index: int,
    slide_count: int,
    first_agenda_index: Optional[int],
    args: argparse.Namespace,
) -> None:
    if not args.check_agenda_page:
        return

    expected_index = min(max(1, args.agenda_page_index), slide_count)
    if first_agenda_index is None:
        if args.require_agenda_page and slide_index == expected_index:
            add_issue(
                issues,
                "FAIL",
                slide_index,
                "AGENDA_PAGE_MISSING",
                f"Deck must include an agenda/table-of-contents page at slide {args.agenda_page_index} "
                "immediately after the mandatory guest-introduction page.",
            )
        return

    if slide_index != first_agenda_index:
        return
    if first_agenda_index == args.agenda_page_index or args.allow_agenda_after_content:
        return

    add_issue(
        issues,
        "FAIL",
        slide_index,
        "AGENDA_PAGE_ORDER",
        f"First agenda/table-of-contents page is slide {first_agenda_index}; it must be slide "
        f"{args.agenda_page_index} immediately after the mandatory guest-introduction page.",
    )


def check_agenda_slide_content(
    issues: list[dict],
    slide_index: int,
    is_agenda: bool,
    text_items: list[tuple],
    slide_h: int,
    args: argparse.Namespace,
) -> None:
    if not args.check_agenda_page or not is_agenda:
        return

    for shape, text in text_items:
        compact = compact_text(text)
        if not compact:
            continue
        normalized = normalize_text(compact)
        forbidden_pattern = next(
            (
                pattern
                for pattern in args.agenda_forbidden_title_pattern
                if normalize_text(pattern) and normalize_text(pattern) in normalized
            ),
            None,
        )
        if forbidden_pattern:
            add_issue(
                issues,
                "FAIL",
                slide_index,
                "AGENDA_FORBIDDEN_TITLE_TEXT",
                f"{shape_label(shape, text)} appears on an agenda page but matches narrative/body-title "
                f"pattern {forbidden_pattern!r}; move it to the first body slide.",
            )
            continue

        if args.allow_agenda_narrative_header or is_agenda_marker_text(compact):
            continue
        rect = shape_rect(shape)
        is_top_header = rect.top / max(1, slide_h) <= args.agenda_header_bottom_ratio
        if is_top_header and weighted_text_len(compact) > args.agenda_header_max_weighted_len:
            add_issue(
                issues,
                "FAIL",
                slide_index,
                "AGENDA_NARRATIVE_HEADER",
                f"{shape_label(shape, text)} is a long top narrative/header on an agenda page; "
                "keep thesis/problem statements on the first body slide.",
            )


def check_pptx(path: Path, args: argparse.Namespace) -> dict:
    prs = Presentation(path)
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    layout_scale = resolve_layout_scale(slide_w, slide_h, args)
    page = Rect(0, 0, slide_w, slide_h)
    margin = inch_to_emu(scaled(args.margin_in, layout_scale))
    footer = Rect(
        0,
        slide_h - inch_to_emu(scaled(args.footer_in, layout_scale)),
        slide_w,
        inch_to_emu(scaled(args.footer_in, layout_scale)),
    )
    logo = Rect(
        slide_w - inch_to_emu(scaled(args.logo_width_in, layout_scale)) - margin,
        margin,
        inch_to_emu(scaled(args.logo_width_in, layout_scale)),
        inch_to_emu(scaled(args.logo_height_in, layout_scale)),
    )
    default_font_pt = scaled(args.default_font_pt, layout_scale)
    warn_min_font_pt = scaled(args.warn_min_font_pt, layout_scale)
    hard_min_font_pt = scaled(args.hard_min_font_pt, layout_scale)
    large_title_pt = scaled(args.large_title_pt, layout_scale)
    issues: list[dict] = []
    slide_statuses = []
    slide_items_cache = []
    agenda_indices: list[int] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        text_items = []
        for shape in walk_shapes(slide.shapes):
            text = shape_text(shape)
            if text:
                text_items.append((shape, text))
        slide_items_cache.append((slide, text_items))
        if is_agenda_slide(slide, [text for _, text in text_items]):
            agenda_indices.append(slide_index)

    first_agenda_index = agenda_indices[0] if agenda_indices else None

    if len(prs.slides) < 2 and not args.skip_mandatory_page_check:
        add_issue(
            issues,
            "FAIL",
            1,
            "MANDATORY_GUEST_PAGE_MISSING",
            "Deck must include a mandatory blank guest-introduction page immediately after the cover.",
        )

    for slide_index, (slide, text_items) in enumerate(slide_items_cache, start=1):
        text_shapes = []
        slide_issues_start = len(issues)
        slide_is_agenda = slide_index in agenda_indices
        has_project_info_style = (
            args.check_project_info_style
            and any(normalize_text(args.project_info_title) in normalize_text(text) for _, text in text_items)
        )
        project_info_item_count = 0

        check_mandatory_slide_content(
            issues,
            slide_index,
            len(prs.slides),
            [text for _, text in text_items],
            args,
        )
        check_key_page_solid_panel(issues, slide_index, len(prs.slides), slide, slide_w, slide_h, args)
        check_agenda_page_order(issues, slide_index, len(prs.slides), first_agenda_index, args)
        check_agenda_slide_content(issues, slide_index, slide_is_agenda, text_items, slide_h, args)

        for shape, text in text_items:

            rect = shape_rect(shape)
            label = shape_label(shape, text)
            text_shapes.append((shape, text, rect))

            check_cover_text_safe_zone(issues, slide_index, shape, text, rect, slide_w, args)

            if (
                rect.left < -args.bounds_tolerance_emu
                or rect.top < -args.bounds_tolerance_emu
                or rect.right > page.right + args.bounds_tolerance_emu
                or rect.bottom > page.bottom + args.bounds_tolerance_emu
            ):
                add_issue(
                    issues,
                    "FAIL",
                    slide_index,
                    "TEXT_OUT_OF_SLIDE",
                    f"{label} extends outside slide bounds at "
                    f"x={emu_to_in(rect.left)}, y={emu_to_in(rect.top)}, "
                    f"w={emu_to_in(rect.width)}, h={emu_to_in(rect.height)}.",
                )

            if intersects(rect, footer):
                level = "FAIL" if args.fail_on_keepout else "WARN"
                add_issue(issues, level, slide_index, "FOOTER_KEEP_OUT", f"{label} intersects footer keep-out area.")

            if intersects(rect, logo):
                level = "FAIL" if args.fail_on_keepout else "WARN"
                add_issue(issues, level, slide_index, "LOGO_KEEP_OUT", f"{label} intersects top-right logo keep-out area.")

            sizes = font_sizes(shape, default_font_pt)
            names = explicit_font_names(shape)
            if slide_index in (1, len(prs.slides)) and not names and not args.allow_theme_fonts_on_key_pages:
                add_issue(
                    issues,
                    "WARN",
                    slide_index,
                    "KEY_PAGE_FONT_NOT_EXPLICIT",
                    f"{label} does not set an explicit font; key cover/closing text may reflow across renderers.",
                )
            min_size = min(sizes)
            max_size = max(sizes)
            if min_size < hard_min_font_pt:
                add_issue(
                    issues,
                    "FAIL",
                    slide_index,
                    "FONT_TOO_SMALL",
                    f"{label} uses {min_size:.1f}pt text, below the scaled "
                    f"{hard_min_font_pt:.1f}pt hard minimum.",
                )
            elif min_size < warn_min_font_pt:
                add_issue(
                    issues,
                    "WARN",
                    slide_index,
                    "FONT_SMALL",
                    f"{label} uses {min_size:.1f}pt text, below the scaled "
                    f"{warn_min_font_pt:.1f}pt warning threshold.",
                )

            estimated_lines = estimate_lines(text, rect.width, max_size)
            required_height_pt = estimated_lines * max_size * args.line_height
            available_height_pt = max(1.0, emu_to_pt(rect.height) - vertical_text_padding_pt(shape, layout_scale))
            if required_height_pt > available_height_pt * args.overflow_tolerance:
                add_issue(
                    issues,
                    "FAIL",
                    slide_index,
                    "TEXT_MAY_OVERFLOW",
                    f"{label} likely needs {required_height_pt:.1f}pt height "
                    f"for about {estimated_lines} lines, available {available_height_pt:.1f}pt.",
                )

            if max_size >= large_title_pt and estimated_lines > args.max_large_title_lines:
                add_issue(
                    issues,
                    "WARN",
                    slide_index,
                    "LARGE_TITLE_WRAP",
                    f"{label} has large text ({max_size:.1f}pt) estimated at {estimated_lines} lines.",
                )

            if (
                args.check_short_info_wrap
                and estimated_lines > 1
                and likely_single_line_info(text, args.short_info_max_weighted_len)
            ):
                level = "FAIL" if args.fail_on_short_info_wrap else "WARN"
                add_issue(
                    issues,
                    level,
                    slide_index,
                    "SHORT_INFO_UNEXPECTED_WRAP",
                    f"{label} is a short key-value/info item estimated at {estimated_lines} lines; "
                    "widen the text box, reduce wording, or explicitly confirm the wrap.",
                )

            if has_project_info_style and is_project_info_item(text):
                project_info_item_count += 1
                level = "FAIL" if args.fail_on_project_info_style else "WARN"
                if estimated_lines > args.project_info_max_item_lines:
                    add_issue(
                        issues,
                        level,
                        slide_index,
                        "PROJECT_INFO_ITEM_WRAP",
                        f"{label} is in a project-info block and is estimated at "
                        f"{estimated_lines} lines; keep project-info facts on one line or log an exception.",
                    )
                min_item_width = inch_to_emu(scaled(args.project_info_min_item_width_in, layout_scale))
                if rect.width < min_item_width:
                    add_issue(
                        issues,
                        level,
                        slide_index,
                        "PROJECT_INFO_ITEM_TOO_NARROW",
                        f"{label} width is {emu_to_in(rect.width)} in, below the scaled "
                        f"{emu_to_in(min_item_width)} in minimum for project-info rows.",
                    )
                if rect.left < int(slide_w * args.project_info_min_left_ratio):
                    add_issue(
                        issues,
                        level,
                        slide_index,
                        "PROJECT_INFO_ITEM_NOT_RIGHT_COLUMN",
                        f"{label} starts at x={emu_to_in(rect.left)} in; project-info facts "
                        "should sit in the right-side fact column for this style.",
                    )
                if args.require_project_info_word_wrap_off and text_word_wrap(shape) is not False:
                    add_issue(
                        issues,
                        level,
                        slide_index,
                        "PROJECT_INFO_ITEM_WORD_WRAP_ENABLED",
                        f"{label} does not explicitly disable word wrap; project-info rows may reflow.",
                    )

        if has_project_info_style and project_info_item_count < args.project_info_min_items:
            level = "FAIL" if args.fail_on_project_info_style else "WARN"
            add_issue(
                issues,
                level,
                slide_index,
                "PROJECT_INFO_TOO_FEW_ITEMS",
                f"Project-info style slide contains {project_info_item_count} recognized fact rows; "
                f"expected at least {args.project_info_min_items}.",
            )

        for i, (_, _, rect_a) in enumerate(text_shapes):
            for shape_b, _, rect_b in text_shapes[i + 1 :]:
                area = intersection_area(rect_a, rect_b)
                if not area:
                    continue
                overlap_ratio = area / max(1, min(rect_a.area, rect_b.area))
                if overlap_ratio >= args.overlap_ratio:
                    label_b = shape_label(shape_b, shape_text(shape_b))
                    add_issue(
                        issues,
                        "FAIL",
                        slide_index,
                        "TEXT_OVERLAP",
                        f"Text shapes overlap by {overlap_ratio:.0%}; later shape is {label_b}.",
                    )

        slide_issues = issues[slide_issues_start:]
        status = "PASS"
        if any(issue["level"] == "FAIL" for issue in slide_issues):
            status = "FAIL"
        elif slide_issues:
            status = "WARN"
        slide_statuses.append({"slide": slide_index, "status": status})

    deck_status = "PASS"
    if any(issue["level"] == "FAIL" for issue in issues):
        deck_status = "FAIL"
    elif issues:
        deck_status = "WARN"

    return {
        "file": str(path),
        "status": deck_status,
        "slide_count": len(prs.slides),
        "canvas_in": {"width": emu_to_in(slide_w), "height": emu_to_in(slide_h)},
        "layout_scale": round(layout_scale, 3),
        "slide_statuses": slide_statuses,
        "issues": issues,
    }


def print_text_report(result: dict) -> None:
    print(
        f"{result['status']} {result['file']} "
        f"({result['slide_count']} slides, layout_scale={result['layout_scale']})"
    )
    for item in result["slide_statuses"]:
        print(f"- slide {item['slide']}: {item['status']}")
    for issue in result["issues"]:
        print(f"[{issue['level']}] slide {issue['slide']} {issue['code']}: {issue['message']}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Check generated PPTX files for layout safety issues.")
    parser.add_argument("pptx", nargs="+", type=Path)
    parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    parser.add_argument("--allow-fail", action="store_true", help="Always exit 0, even when FAIL issues are found.")
    parser.add_argument("--fail-on-keepout", action="store_true", help="Treat logo/footer keep-out intersections as FAIL.")
    parser.add_argument("--layout-scale", type=float, default=None, help="Override automatic 26.667x15 based scaling.")
    parser.add_argument("--no-auto-scale", action="store_true", help="Disable automatic half-size threshold scaling.")
    parser.add_argument("--skip-mandatory-page-check", action="store_true", help="Skip 嘉宾介绍/感谢聆听 hard checks.")
    parser.add_argument("--allow-guest-page-content", action="store_true", help="Allow content on the mandatory 嘉宾介绍 page.")
    parser.add_argument("--allow-custom-closing", action="store_true", help="Allow a final slide without the default 感谢聆听 copy.")
    parser.add_argument("--guest-page-title", default="嘉宾介绍")
    parser.add_argument("--closing-text", default="感谢聆听")
    parser.add_argument("--disable-cover-safe-zone", action="store_true", help="Skip the cover text-safe-zone check.")
    parser.add_argument("--cover-text-right-ratio", type=float, default=0.33)
    parser.add_argument("--cover-left-column-center-ratio", type=float, default=0.43)
    parser.add_argument("--allow-theme-fonts-on-key-pages", action="store_true", help="Do not warn when cover/closing text uses theme fonts.")
    parser.add_argument("--check-key-page-solid-panel", action=argparse.BooleanOptionalAction, default=True)
    parser.add_argument("--allow-native-cover-layout", action="store_true", help="Allow title page without the default centered solid panel after explicit user approval.")
    parser.add_argument("--allow-native-closing-layout", action="store_true", help="Allow final page without the default centered solid panel after explicit user approval.")
    parser.add_argument("--key-panel-color", default="061C4C", help="Expected deep-blue solid panel color as 6-digit RGB hex.")
    parser.add_argument("--key-panel-color-tolerance", type=int, default=48)
    parser.add_argument("--key-panel-min-width-ratio", type=float, default=0.30)
    parser.add_argument("--key-panel-max-width-ratio", type=float, default=0.86)
    parser.add_argument("--key-panel-min-height-ratio", type=float, default=0.12)
    parser.add_argument("--key-panel-max-height-ratio", type=float, default=0.58)
    parser.add_argument("--key-panel-min-area-ratio", type=float, default=0.035)
    parser.add_argument("--key-panel-center-x-tolerance", type=float, default=0.18)
    parser.add_argument("--key-panel-center-y-tolerance", type=float, default=0.24)
    parser.add_argument("--check-agenda-page", action=argparse.BooleanOptionalAction, default=True)
    parser.add_argument("--require-agenda-page", action="store_true", help="Require an agenda/table-of-contents page at --agenda-page-index.")
    parser.add_argument("--allow-agenda-after-content", action="store_true", help="Allow the first agenda page to appear after body content when explicitly approved.")
    parser.add_argument("--allow-agenda-narrative-header", action="store_true", help="Allow long narrative headers on agenda pages after explicit approval.")
    parser.add_argument("--agenda-page-index", type=int, default=3)
    parser.add_argument("--agenda-header-bottom-ratio", type=float, default=0.26)
    parser.add_argument("--agenda-header-max-weighted-len", type=float, default=18.0)
    parser.add_argument("--agenda-forbidden-title-pattern", action="append", default=list(DEFAULT_AGENDA_FORBIDDEN_PATTERNS))
    parser.add_argument("--margin-in", type=float, default=0.45)
    parser.add_argument("--footer-in", type=float, default=0.75)
    parser.add_argument("--logo-width-in", type=float, default=3.2)
    parser.add_argument("--logo-height-in", type=float, default=1.0)
    parser.add_argument("--default-font-pt", type=float, default=20.0)
    parser.add_argument("--warn-min-font-pt", type=float, default=14.0)
    parser.add_argument("--hard-min-font-pt", type=float, default=12.0)
    parser.add_argument("--large-title-pt", type=float, default=38.0)
    parser.add_argument("--max-large-title-lines", type=int, default=2)
    parser.add_argument("--check-short-info-wrap", action=argparse.BooleanOptionalAction, default=True)
    parser.add_argument("--fail-on-short-info-wrap", action="store_true", help="Treat short key-value/info wraps as FAIL.")
    parser.add_argument("--short-info-max-weighted-len", type=float, default=38.0)
    parser.add_argument("--check-project-info-style", action=argparse.BooleanOptionalAction, default=True)
    parser.add_argument("--fail-on-project-info-style", action="store_true", help="Treat project-info style issues as FAIL.")
    parser.add_argument("--project-info-title", default="项目基本信息")
    parser.add_argument("--project-info-min-items", type=int, default=3)
    parser.add_argument("--project-info-max-item-lines", type=int, default=1)
    parser.add_argument("--project-info-min-item-width-in", type=float, default=7.4)
    parser.add_argument("--project-info-min-left-ratio", type=float, default=0.52)
    parser.add_argument("--require-project-info-word-wrap-off", action=argparse.BooleanOptionalAction, default=True)
    parser.add_argument("--line-height", type=float, default=1.2)
    parser.add_argument("--overflow-tolerance", type=float, default=1.0)
    parser.add_argument("--overlap-ratio", type=float, default=0.08)
    parser.add_argument("--bounds-tolerance-emu", type=int, default=1000)
    args = parser.parse_args()

    results = [check_pptx(path, args) for path in args.pptx]
    if args.json:
        print(json.dumps(results, ensure_ascii=False, indent=2))
    else:
        for index, result in enumerate(results):
            if index:
                print()
            print_text_report(result)

    has_fail = any(result["status"] == "FAIL" for result in results)
    if has_fail and not args.allow_fail:
        raise SystemExit(1)


if __name__ == "__main__":
    main()

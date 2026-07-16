#!/usr/bin/env python3
"""Read-only PPTX style inspector for TGO/GTLC conversion work."""

from __future__ import annotations

import argparse
import json
import sys
from collections import Counter
from io import BytesIO
from pathlib import Path
from zipfile import ZipFile

etree = None
Image = None
Presentation = None
MSO_SHAPE_TYPE = None


def load_pptx_style_dependencies() -> list[str]:
    """Load optional inspection dependencies after CLI parsing."""
    global etree, Image, Presentation, MSO_SHAPE_TYPE
    missing: list[str] = []
    try:
        from lxml import etree as _etree
    except ImportError:
        missing.append("lxml")
    else:
        etree = _etree
    try:
        from PIL import Image as _Image
    except ImportError:
        missing.append("Pillow")
    else:
        Image = _Image
    try:
        from pptx import Presentation as _Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO_SHAPE_TYPE
    except ImportError:
        missing.append("python-pptx")
    else:
        Presentation = _Presentation
        MSO_SHAPE_TYPE = _MSO_SHAPE_TYPE
    return missing

EMU_PER_IN = 914400
EMU_PER_PT = 12700
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def emu_to_in(value):
    return round(value / EMU_PER_IN, 3) if value is not None else None


def emu_to_pt(value):
    return round(value / EMU_PER_PT, 1) if value is not None else None


def color_value(node):
    if node is None:
        return None
    srgb = node.find(".//a:srgbClr", NS)
    scheme = node.find(".//a:schemeClr", NS)
    sys = node.find(".//a:sysClr", NS)
    if srgb is not None:
        return "#" + srgb.get("val")
    if scheme is not None:
        return "scheme:" + scheme.get("val")
    if sys is not None:
        return "sys:" + sys.get("val") + " last=#" + str(sys.get("lastClr"))
    return None


def color_from_python_pptx(color):
    try:
        if color.rgb:
            return "#" + str(color.rgb)
    except Exception:
        pass
    try:
        if color.theme_color is not None:
            return "theme:" + str(color.theme_color)
    except Exception:
        pass
    return None


def shape_text(shape):
    try:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            return " ".join(shape.text.split())
    except Exception:
        return ""
    return ""


def walk_shapes(shapes):
    for shape in shapes:
        yield shape
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk_shapes(shape.shapes)
        except Exception:
            pass


def collect_layouts(prs):
    layouts = []
    for layout in prs.slide_layouts:
        placeholders = []
        pictures = []
        text_boxes = []
        fills = Counter()
        for shape in walk_shapes(layout.shapes):
            item = {
                "name": shape.name,
                "x": emu_to_in(getattr(shape, "left", None)),
                "y": emu_to_in(getattr(shape, "top", None)),
                "w": emu_to_in(getattr(shape, "width", None)),
                "h": emu_to_in(getattr(shape, "height", None)),
            }
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures.append(item)
            text = shape_text(shape)
            if text:
                item["text"] = text[:120]
                text_boxes.append(item)
            try:
                item["placeholder_type"] = str(shape.placeholder_format.type)
                placeholders.append(item)
            except Exception:
                pass
            try:
                fill = color_from_python_pptx(shape.fill.fore_color)
                if fill:
                    fills[fill] += 1
            except Exception:
                pass
        layouts.append(
            {
                "name": layout.name,
                "placeholders": placeholders,
                "pictures": pictures,
                "text_boxes": text_boxes,
                "fills": dict(fills),
            }
        )
    return layouts


def collect_theme_and_text_styles(path):
    sizes = Counter()
    fonts = Counter()
    colors = Counter()
    theme_colors = []
    with ZipFile(path) as zf:
        xml_files = [
            name
            for name in zf.namelist()
            if name.startswith(("ppt/theme/theme", "ppt/slideMasters/slideMaster", "ppt/slideLayouts/slideLayout"))
            and name.endswith(".xml")
        ]
        for name in xml_files:
            xml = etree.fromstring(zf.read(name))
            for child in xml.xpath(".//a:clrScheme/*", namespaces=NS):
                tag = etree.QName(child).localname
                theme_colors.append({"role": tag, "value": color_value(child)})
            for node in xml.xpath(".//a:rPr|.//a:defRPr|.//a:endParaRPr", namespaces=NS):
                size = node.get("sz")
                if size:
                    sizes[round(int(size) / 100, 1)] += 1
                latin = node.find("./a:latin", NS)
                if latin is not None and latin.get("typeface"):
                    fonts[latin.get("typeface")] += 1
                color = color_value(node)
                if color:
                    colors[color] += 1
    return {
        "theme_colors": theme_colors,
        "text_sizes_pt": dict(sizes.most_common()),
        "fonts": dict(fonts.most_common()),
        "text_colors": dict(colors.most_common()),
    }


def collect_media(path):
    media = []
    with ZipFile(path) as zf:
        for name in sorted(n for n in zf.namelist() if n.startswith("ppt/media/") and not n.endswith("/")):
            data = zf.read(name)
            try:
                image = Image.open(BytesIO(data)).convert("RGB")
                media.append({"name": name, "width": image.width, "height": image.height, "bytes": len(data)})
            except Exception:
                media.append({"name": name, "bytes": len(data)})
    return media


def inspect(path):
    prs = Presentation(path)
    result = {
        "file": str(path),
        "canvas": {
            "width_in": emu_to_in(prs.slide_width),
            "height_in": emu_to_in(prs.slide_height),
        },
        "slide_count": len(prs.slides),
        "master_count": len(prs.slide_masters),
        "layouts": collect_layouts(prs),
        "media": collect_media(path),
    }
    result.update(collect_theme_and_text_styles(path))
    return result


def main():
    parser = argparse.ArgumentParser(description="Inspect PPTX dimensions, layouts, theme, text styles, and media.")
    parser.add_argument("pptx", nargs="+", type=Path)
    args = parser.parse_args()
    missing = load_pptx_style_dependencies()
    if missing:
        distributions = ", ".join(dict.fromkeys(missing))
        print(
            f"E_DEPENDENCY_MISSING: required distributions are unavailable: {distributions}; "
            "install them and retry.",
            file=sys.stderr,
        )
        raise SystemExit(2)
    print(json.dumps([inspect(path) for path in args.pptx], ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
